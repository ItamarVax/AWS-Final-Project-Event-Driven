#!/usr/bin/env python3
"""Render the HTML slides to images via headless Chrome and assemble a
16:9 PPTX with full-bleed slide images + speaker notes (from the outline).

Owner split: slides 1-6 = PARTNER (Maor), slides 7-16 = YOU (Itamar).
Each note carries an owner+time tag and 3 likely audience Q&A."""
import json, os, subprocess, sys
from pptx import Presentation
from pptx.util import Inches

ROOT = os.path.dirname(os.path.abspath(__file__))
SLIDES_DIR = os.path.join(ROOT, "slides")
BUILD_DIR = os.path.join(ROOT, "build")
OUT_PPTX = os.path.join(ROOT, "AWS-Bedrock-Presentation-Notes.pptx")
CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"

# Speaker notes, keyed by slide filename. Sourced from bedrock-presentation-outline.md.
NOTES = {
"01-title.html": """[PARTNER — Maor · 0:30] Introduce the team — Itamar Vax & Maor Aharon. One-sentence promise: "By the end you'll know what Bedrock is, how it works under the hood, and you'll watch it classify real PII data live." Flag the structure: Part 1 = what it is + how it works; Part 2 = in practice + live demo.

Likely questions:
Q: Why Bedrock and not just call OpenAI directly? — A: Everything we run is on AWS; Bedrock gives many models behind our existing IAM auth, with no separate keys or accounts. (We expand on this in a moment.)
Q: Is this a real project or a class exercise? — A: It's a real production pipeline at BigID; we use it as the running example throughout.
Q: What will the live demo show? — A: A real PII finding classified by a model through Bedrock, returning structured JSON live.""",

"02-agenda.html": """[PARTNER — Maor · 1:00] Open with the hook: "Every LLM call in a production system we run at BigID goes through Bedrock." The hook earns attention — this isn't a toy, it's a system we operate. Walk the roadmap: What → How it works → Alternatives & trade-offs → Real use case → Live demo. Flag that the second half is hands-on.

Likely questions:
Q: How long is the talk? — A: About 24 minutes, with a live demo in the second half.
Q: Do we need AWS or ML background to follow? — A: No — the first half is concepts; the code and demo come later.
Q: What's BigID? — A: A data-security company; the pipeline we show classifies PII (personal-data) findings.""",

"03-what-is-bedrock.html": """[PARTNER — Maor · 1:30] Core definition: one fully-managed AWS API fronting many foundation models (Claude, Nova/Titan, Llama, Mistral, Cohere). Mental model: "Bedrock is to LLMs what S3 is to storage — a managed AWS service you call, not infrastructure you run." Emphasize serverless (course tie-in): you never provision a GPU; pay per token; scales to zero.

Likely questions:
Q: What's a "foundation model"? — A: A large, pre-trained model (like Claude or Llama) you consume via API rather than training yourself.
Q: Is Bedrock itself a model? — A: No — it's a managed gateway/API in front of many vendors' models.
Q: What does "scales to zero" mean? — A: No always-on servers; idle time costs nothing, you pay only for the tokens you actually use.""",

"04-problem-it-solves.html": """[PARTNER — Maor · 2:00] Lead with the pain: wanting Claude AND Llama AND Nova without Bedrock means three SDKs, three key stores, three request formats. Bedrock collapses that into one boto3 client + one Converse shape. The AWS-native win is IAM auth — credentials are your existing role, nothing extra to store, rotate, or leak. The take: Bedrock's win isn't "another LLM service," it's collapsing three into one IAM-native API.

Likely questions:
Q: Can't you just use one vendor's SDK? — A: You can, but using several models then means several SDKs, key stores, and formats — Bedrock unifies them.
Q: What is IAM auth and why is it better than API keys? — A: IAM is AWS's built-in identity system; you reuse your existing role, so there are no separate keys to store, rotate, or leak.
Q: Does this convenience cost extra? — A: You pay per token, roughly at vendor list price; the win is operational simplicity, not a markup.""",

"05-where-it-fits.html": """[PARTNER — Maor · 1:00] Two planes: control plane (bedrock) = manage/list models, quotas, access — admin calls you touch rarely. Runtime plane (bedrock-runtime) = call models via converse()/converse_stream(); token-billed, latency-sensitive — the hot path you live in. Regional, IAM-scoped, integrates with Lambda / Step Functions / CloudWatch / VPC endpoints. You only touch bedrock-runtime for inference — and that's one line of code (next slide).

Likely questions:
Q: What's the difference between "bedrock" and "bedrock-runtime"? — A: bedrock (control plane) manages and lists models and quotas; bedrock-runtime (data plane) actually runs inference. You mostly use runtime.
Q: Is Bedrock regional or global? — A: Regional — you call a regional endpoint, and inference profiles can route across regions.
Q: How does it fit a serverless app? — A: It's a managed API you call from Lambda/Step Functions; nothing to run, and it plugs into CloudWatch and IAM.""",

"06-architecture.html": """[PARTNER — Maor · 2:00] Walk the arrow: the app never holds the model; it sends an IAM-signed HTTPS request to a regional endpoint; AWS routes to the model by inference profile; telemetry flows to CloudWatch. Point out the Config (600s read timeout, 6 retries) — production LLM calls are slow and need resilience. Four components per call: (1) runtime client, (2) model/inference-profile ID, (3) request = system + messages + inferenceConfig, (4) response = content + token usage. (Stay at the picture level — Itamar takes the code-deep slides next.)

Likely questions:
Q: What's the difference between boto3 and aioboto3? — A: boto3 is AWS's Python SDK; aioboto3 is its async version, which lets us run many classifications concurrently.
Q: Why a 600-second timeout and 6 retries? — A: LLM calls can be slow and occasionally throttle; a long timeout plus retries makes production calls resilient.
Q: What's an "inference profile"? — A: An ID (e.g. us.anthropic.…) that routes the request across regions for throughput and availability — Itamar shows this on the next slide.""",

"07-converse-api.html": """[YOU — Itamar · 2:00] Headline: "This exact dict works for Claude, Nova, Llama, Mistral — I only change modelId." That's the whole value prop in one block. Converse unifies the request shape across vendors. Inference profiles: the us. prefix (us.anthropic.claude-sonnet-4-6) load-balances across US regions for throughput and availability, with no code change. Plant the seed: "remember this — it pays off as a lesson at the end."

Likely questions:
Q: Does Converse work identically for every model? — A: The request shape is identical (system / messages / inferenceConfig); a few advanced features vary by model.
Q: What exactly changes when you swap models? — A: Just the modelId string — everything else in the request stays the same.
Q: Was there an API before Converse? — A: Yes, InvokeModel with per-vendor JSON bodies; Converse replaced that with one unified shape across vendors.""",

"08-tool-use.html": """[YOU — Itamar · 0:45] Standout reliability feature. Problem: LLMs return prose; production code needs a typed dict, and regex-on-prose is how demos pass and pipelines fail at 3am. Fix: define a Pydantic model (e.g. {score, reasoning}); Bedrock exposes it as a tool; toolChoice forces the call — output is valid, typed JSON by construction. We never string-parse prose. The difference between a demo and a production system.

Likely questions:
Q: How does tool-use force valid JSON? — A: You expose a schema as a "tool" and set toolChoice to that tool, so the model must respond by calling it with structured input.
Q: What if the model still returns malformed output? — A: With forced tool-use the response block is structured by construction; we read toolUse.input directly, so there's nothing to parse.
Q: Is this the same as OpenAI's function calling? — A: Conceptually yes — tool-use is Bedrock's equivalent, unified across vendors through Converse.""",

"09-cost-quality-levers.html": """[YOU — Itamar · 0:45] Two more real levers. Prompt caching = money: our fixed system prompt (the classifier "brief") is reused across thousands of findings — pay full price once, cache reads bill at 10% of the input rate. Extended thinking = quality: let the model reason harder on genuinely hard calls. Honest nuance: forced tool-use + thinking can't combine, so we relax to toolChoice: auto and the parser handles the fallback.

Likely questions:
Q: How much does prompt caching actually save? — A: Cache reads bill at ~10% of the input-token rate, so a reused system prompt is ~10x cheaper after the first call.
Q: What is extended thinking? — A: The model spends extra tokens reasoning before answering — you pay more but decide better on genuinely hard cases.
Q: Why can't forced tool-use and thinking combine? — A: Those model modes are mutually exclusive, so we relax to toolChoice: auto and let the parser handle the fallback.""",

"10-alternatives.html": """[YOU — Itamar · 2:00] Position Bedrock in the middle: more managed than SageMaker (self-host = you own instances/scaling/ops/GPU bill), more flexible than Amazon Q (opinionated assistant, not raw model access); Lambda + external API gives the broadest menu but you manage keys and lose IAM-native auth. Concrete contrast from our repo: we DO use OpenAI direct + OpenRouter for two sub-tasks — and each needed its own wrapper module, unlike Bedrock's single Converse call. That sets up the slide-15 lesson.

Likely questions:
Q: When would you pick SageMaker over Bedrock? — A: When you need full control — custom or self-hosted models, your own fine-tuning and endpoints — and accept managing the infrastructure.
Q: What is Amazon Q? — A: A higher-level, prebuilt AWS assistant; convenient for workflows but not raw model access like Bedrock.
Q: Why do you still call OpenAI directly in the repo? — A: Two sub-tasks needed models/features outside our Bedrock setup; each needed its own wrapper — which is exactly the proof of Bedrock's single-API value.""",

"11-advantages-limitations.html": """[YOU — Itamar · 2:00] Keep advantages crisp (one API/many models, IAM auth, managed/serverless, caching+tool-use+thinking built in, per-model cost tracking). Spend the time on limitations because they're real and earned — each maps to actual code we wrote: region-anchored inference profiles + per-account/region model activation; Llama/Mistral/Cohere reject cachePoint (we added BEDROCK_CACHE_SYSTEM=false); CountTokens needs the single-region ID (we strip the us. prefix); forced tool-use + thinking can't combine; quotas/throughput/regional availability vary; and Learner Lab may not expose all models — hence the recorded demo fallback.

Likely questions:
Q: What's the biggest limitation in practice? — A: Per-account/region model activation plus feature gaps (some models reject caching); you plan around model+region pairings.
Q: Does every model support every feature? — A: No — Llama/Mistral/Cohere rejected cachePoint, so we added a toggle to disable caching for them.
Q: Isn't there vendor lock-in? — A: Somewhat, but our provider seam (core/llm.py) makes swapping providers a one-elif change.""",

"12-use-case-pipeline.html": """[YOU — Itamar · 2:00] This is the "real-world use case" deliverable. The problem: regex flags candidate PII with many false positives; we need an LLM to judge true vs false positive per PII type, guided by a classifier brief. Every box in the loop (spec → prompt → dataset → judge → validate, iterate) is a Bedrock call. The system automatically writes and refines the prompt a downstream validator uses — meta, but concrete. Don't go deep on PII semantics; keep focus on "multi-agent system, Bedrock serves every model."

Likely questions:
Q: What does the pipeline actually do? — A: Regex flags candidate PII; LLMs then judge true vs false positive per type and automatically refine the classifier prompt.
Q: Why multiple agents instead of one big prompt? — A: Each step (spec, generate, judge, validate) is a focused task; separating them improves reliability and lets us pick the right model per step.
Q: Are you running this on real personal data? — A: For the demo we use synthetic data; nothing sensitive is shown.""",

"13-model-mix-cost.html": """[YOU — Itamar · 2:00] The most persuasive slide — it turns "multi-model" from a feature into dollars. Tell the story: you don't pay Opus prices to label 10,000 cheap findings; you use Nova 2 Lite / Haiku for volume and reserve Opus for the hard judge step. Nova 2 Lite input is ~83x cheaper than Opus; Haiku ~20x cheaper; caching bills reused prompts at 10% on reads. Because of Converse, choosing the cost-optimal model per task is a one-string change, not an engineering project.

Likely questions:
Q: Why not just use the best model everywhere? — A: Cost — Opus is ~83x pricier on input than Nova Lite; you reserve it for the hard judge step only.
Q: How hard is it to switch models per task? — A: One string change (the modelId / inference-profile), thanks to Converse — no new code.
Q: Are these prices current? — A: They're representative Bedrock list prices per million tokens; exact figures vary by region and over time.""",

"14-live-demo.html": """[YOU — Itamar · ~4:00 — TIME BUFFER] Live demo. Goal: real input → Bedrock → real structured output that supports the explanation, not just clicking the console. Run validate_finding_async — the innermost call: one PII finding + classifier brief → Bedrock Converse + forced tool-use → {"score": 1|2|3, "reasoning": "..."}. Watch for: the JSON shape and the token/cost log line (ties back to slides 9 & 13). Use synthetic PII so nothing sensitive is on screen. Have the recorded fallback ready if live Bedrock access fails. Trim or extend here to land on time.

Likely questions:
Q: What are we actually looking at? — A: validate_finding_async — one PII finding plus the classifier brief going through Bedrock Converse with forced tool-use, returning {score, reasoning}.
Q: What do the score values mean? — A: 1/2/3 encode the validator's judgment (confidence that a finding is a true vs false positive) per our schema.
Q: What happens if the live call fails? — A: We have a recorded fallback of the same run ready to play.""",

"15-lessons-learned.html": """[YOU — Itamar · 1:00] Land the Converse seed planted on slide 7. Five earned lessons: (1) one API, not N wrappers — Converse meant zero per-provider request code for Bedrock models, while our OpenAI/OpenRouter paths each needed a wrapper, which proves the point; (2) tool-use beats prose parsing; (3) cache the system prompt — biggest single-line cost win; (4) build a provider seam (core/llm.py) so adding a provider is a one-elif change; (5) mind the gotchas (region-anchored profiles, CountTokens single-region ID, not all models cache). Earned lessons from building the system — exactly what the assignment asks for.

Likely questions:
Q: What's the single biggest takeaway? — A: One API (Converse) eliminated per-provider code for Bedrock models — swapping a model is a one-string change.
Q: What surprised you most? — A: How much prompt caching saved (10% read rate) on a reused system prompt — the biggest cost win in the project.
Q: What would you do differently? — A: Build the provider seam (core/llm.py) earlier, so adding any provider stays a one-elif change.""",

"16-wrap-qa.html": """[YOU — Itamar · 1:00] Three takeaways: (1) Bedrock = managed, multi-model GenAI on AWS — serverless, IAM-auth, scales to zero; (2) Converse = one request shape for every model, swap with a one-string change; (3) right-model-for-the-job economics — Nova/Haiku for volume, Opus for hard calls. Both presenters field questions. If short on time, recover here; if long, the demo already absorbed it. Close: "Thank you — see the repo for the full code path."

Likely questions:
Q: Bottom line — why Bedrock? — A: Managed, multi-model GenAI on AWS: IAM auth, serverless, swap models with one string, pay per token.
Q: Could we use this outside AWS? — A: The concepts transfer, but the IAM-native auth and AWS integration are the Bedrock-specific wins.
Q: Where can we see the code? — A: It's in our repo — the full path from client creation through to the validator call.""",
}

def render(playlist):
    os.makedirs(BUILD_DIR, exist_ok=True)
    pngs = []
    for name in playlist:
        src = os.path.join(SLIDES_DIR, name)
        out = os.path.join(BUILD_DIR, name.replace(".html", ".png"))
        cmd = [CHROME, "--headless=new", "--disable-gpu", "--hide-scrollbars",
               "--force-device-scale-factor=2", "--window-size=1920,1080",
               "--virtual-time-budget=8000",
               f"--screenshot={out}", f"file://{src}"]
        print("render", name, flush=True)
        subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        if not os.path.exists(out):
            sys.exit(f"FAILED to render {name}")
        pngs.append((name, out))
    return pngs

def build(pngs):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for name, png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)
        note = NOTES.get(name, "")
        slide.notes_slide.notes_text_frame.text = note
    prs.save(OUT_PPTX)
    print("saved", OUT_PPTX, f"({os.path.getsize(OUT_PPTX)//1024} KB)")

if __name__ == "__main__":
    with open(os.path.join(ROOT, "manifest.json")) as f:
        playlist = json.load(f)["playlist"]
    missing = [n for n in playlist if n not in NOTES]
    if missing:
        sys.exit(f"missing notes for: {missing}")
    build(render(playlist))
